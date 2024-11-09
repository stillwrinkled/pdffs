import io
import os
import zipfile
from pdf2docx import Converter as PDF2DOCXConverter
import fitz  # PyMuPDF for converting PDF to images (JPEG/PNG)
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
import camelot
from PyPDF2 import PdfReader
import tempfile


def convert_pdf(file, format):
    output = io.BytesIO()

    # Reset file pointer to the start in case it has been read before
    file.seek(0)

    if format == 'docx':
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
                temp_pdf.write(file.read())
                temp_pdf.close()
                output_file = io.BytesIO()
                converter = PDF2DOCXConverter(temp_pdf.name)
                converter.convert(output_file, start=0, end=None)
                converter.close()
                os.remove(temp_pdf.name)
                output_file.seek(0)

                # Set a filename for the DOCX file
                output_filename = "converted_file.docx"
                return output_file, output_filename

        except Exception as e:
            print(f"Error during PDF to DOCX conversion: {e}")
            raise ValueError("Failed to convert PDF to DOCX.")

    elif format == 'xlsx':
        tables = camelot.read_pdf(file, pages='all')
        wb = Workbook()
        ws = wb.active
        for table in tables:
            for i, row in enumerate(table.df.itertuples(), 1):
                for j, value in enumerate(row[1:], 1):
                    ws.cell(row=i, column=j, value=value)
        wb.save(output)

    elif format == 'pptx':
        prs = Presentation()
        reader = PdfReader(file)
        for page_num, page in enumerate(reader.pages):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            text_frame = text_box.text_frame
            text_frame.text = page.extract_text()
        prs.save(output)

    elif format in ['jpeg', 'png']:
        try:
            # Check if it's a multi-page PDF and adjust the output type
            doc = fitz.open(stream=file.read(), filetype="pdf")
            img_format = format.upper()

            if len(doc) > 1:
                # Create a zip file in memory to store the images for multi-page PDFs
                with zipfile.ZipFile(output, mode='w', compression=zipfile.ZIP_DEFLATED) as img_zip:
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)
                        pix = page.get_pixmap()
                        img_bytes = pix.tobytes(img_format)

                        # Write each image to a separate file in the zip archive
                        img_filename = f'page_{page_num + 1}.{format}'  # E.g., page_1.png or page_1.jpeg
                        img_zip.writestr(img_filename, img_bytes)

                # Set the correct output extension for multi-page PDFs
                output_filename = "converted_pages.zip"
                output.seek(0)
                return output, output_filename

            else:
                # Single-page output, return the image directly
                pix = doc.load_page(0).get_pixmap()
                img_bytes = pix.tobytes(img_format)
                output.write(img_bytes)
                output_filename = f"converted_page.{format}"  # Set the correct image format
                output.seek(0)
                return output, output_filename

        except fitz.FileDataError as e:
            print(f"Error: {e}")
            raise ValueError("Failed to process the PDF file. Please check if it's a valid PDF.")

    output.seek(0)
    return output, None


def html_to_pdf(html_content):
    from xhtml2pdf import pisa
    output = io.BytesIO()
    pisa.CreatePDF(html_content, dest=output)
    output.seek(0)
    return output
