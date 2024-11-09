import io
import os
import zipfile
from pdf2docx import Converter as PDF2DOCXConverter
import fitz  # PyMuPDF for converting PDF to images (JPEG/PNG)
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
import camelot  # For extracting tables from PDF
import tabula  # Alternative for extracting tables using Java
from PyPDF2 import PdfReader
import tempfile
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def convert_pdf(file, format):
    output = io.BytesIO()
    file.seek(0)  # Reset file pointer to the start

    if format == 'docx':
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
                temp_pdf.write(file.read())
                temp_pdf.close()
                output_file = io.BytesIO()
                converter = PDF2DOCXConverter(temp_pdf.name)
                converter.convert(output_file, start=0, end=None)
                converter.close()
                os.remove(temp_pdf.name)
                output_file.seek(0)
                output_filename = "converted_file.docx"
                logger.info("PDF to DOCX conversion successful")
                return output_file, output_filename

        except Exception as e:
            logger.error(f"Error during PDF to DOCX conversion: {e}")
            raise ValueError("Failed to convert PDF to DOCX.")

    elif format == 'xlsx':
        try:
            # Attempt to extract tables using Camelot
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
                temp_pdf.write(file.read())
                temp_pdf.close()

                try:
                    logger.info("Attempting table extraction with Camelot")
                    tables = camelot.read_pdf(temp_pdf.name, pages='all', flavor='stream')

                    if not tables or len(tables) == 0:
                        raise ValueError("No tables found in the PDF.")

                    wb = Workbook()
                    ws = wb.active

                    # Write extracted tables to Excel
                    for table_num, table in enumerate(tables):
                        for i, row in enumerate(table.df.itertuples(), 1):
                            for j, value in enumerate(row[1:], 1):
                                ws.cell(row=i, column=j, value=value)

                    wb.save(output)
                    output_filename = "converted_file.xlsx"
                    output.seek(0)
                    os.remove(temp_pdf.name)
                    logger.info("PDF to Excel conversion successful using Camelot")
                    return output, output_filename

                except Exception as camelot_err:
                    logger.error(f"Camelot failed: {camelot_err}")
                    # If Camelot fails, try using Tabula
                    logger.info("Camelot failed, trying Tabula")
                    try:
                        tables = tabula.read_pdf(temp_pdf.name, pages='all', multiple_tables=True)
                        if not tables or len(tables) == 0:
                            raise ValueError("No tables found using Tabula.")

                        wb = Workbook()
                        ws = wb.active

                        # Write extracted tables to Excel
                        for table in tables:
                            for row in table.itertuples():
                                ws.append(row)

                        wb.save(output)
                        output_filename = "converted_file.xlsx"
                        output.seek(0)
                        os.remove(temp_pdf.name)
                        logger.info("PDF to Excel conversion successful using Tabula")
                        return output, output_filename

                    except Exception as tabula_err:
                        logger.error(f"Tabula failed: {tabula_err}")
                        raise ValueError("Failed to extract tables with both Camelot and Tabula.")

        except Exception as e:
            logger.error(f"Error during PDF to Excel conversion: {e}")
            raise ValueError("Failed to convert PDF to Excel. Please check if the PDF contains valid tables.")

    elif format == 'pptx':
        try:
            prs = Presentation()
            reader = PdfReader(file)
            for page_num, page in enumerate(reader.pages):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                text_frame = text_box.text_frame
                text_frame.text = page.extract_text()
            prs.save(output)
            output_filename = "converted_file.pptx"
            output.seek(0)
            logger.info("PDF to PPTX conversion successful")
            return output, output_filename

        except Exception as e:
            logger.error(f"Error during PDF to PPTX conversion: {e}")
            raise ValueError("Failed to convert PDF to PPTX.")

    elif format in ['jpeg', 'png']:
        try:
            doc = fitz.open(stream=file.read(), filetype="pdf")
            img_format = format.upper()

            if len(doc) > 1:
                with zipfile.ZipFile(output, mode='w', compression=zipfile.ZIP_DEFLATED) as img_zip:
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)
                        pix = page.get_pixmap()
                        img_bytes = pix.tobytes(img_format)
                        img_filename = f'page_{page_num + 1}.{format}'
                        img_zip.writestr(img_filename, img_bytes)

                output_filename = "converted_pages.zip"
                output.seek(0)
                logger.info("PDF to multi-page image (ZIP) conversion successful")
                return output, output_filename

            else:
                pix = doc.load_page(0).get_pixmap()
                img_bytes = pix.tobytes(img_format)
                output.write(img_bytes)
                output_filename = f"converted_page.{format}"
                output.seek(0)
                logger.info("PDF to single image conversion successful")
                return output, output_filename

        except fitz.FileDataError as e:
            logger.error(f"Error during image conversion: {e}")
            raise ValueError("Failed to process the PDF file. Please check if it's a valid PDF.")

    output.seek(0)
    return output, None


def html_to_pdf(html_content):
    from xhtml2pdf import pisa
    output = io.BytesIO()
    try:
        pisa.CreatePDF(html_content, dest=output)
        output.seek(0)
        logger.info("HTML to PDF conversion successful")
    except Exception as e:
        logger.error(f"Error during HTML to PDF conversion: {e}")
        raise ValueError("Failed to convert HTML to PDF.")
    return output
