import io
import os
import zipfile
from pdf2docx import Converter as PDF2DOCXConverter
import fitz  # PyMuPDF for converting PDF to images (JPEG/PNG)
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
import camelot  # For extracting tables from PDF
from PyPDF2 import PdfReader
import tempfile

def convert_pdf(file, format):
    output = io.BytesIO()

    # Reset file pointer to the start in case it has been read before
    file.seek(0)

    if format == 'docx':
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
                temp_pdf.write(file.read())
                temp_pdf.close()
                output_file = io.BytesIO()
                converter = PDF2DOCXConverter(temp_pdf.name)
                converter.convert(output_file, start=0, end=None)
                converter.close()
                os.remove(temp_pdf.name)
                output_file.seek(0)
                output_filename = "converted_file.docx"
                return output_file, output_filename

        except Exception as e:
            print(f"Error during PDF to DOCX conversion: {e}")
            raise ValueError("Failed to convert PDF to DOCX.")

    elif format == 'xlsx':
        try:
            # Using Camelot in 'stream' mode to extract tables from PDF
            tables = camelot.read_pdf(file, pages='all', flavor='stream')

            if not tables or len(tables) == 0:
                raise ValueError("No tables found in the PDF or table extraction failed.")

            wb = Workbook()
            ws = wb.active

            # Loop through all detected tables and write them to Excel
            for table_num, table in enumerate(tables):
                for i, row in enumerate(table.df.itertuples(), 1):
                    for j, value in enumerate(row[1:], 1):
                        ws.cell(row=i, column=j, value=value)

            wb.save(output)
            output_filename = "converted_file.xlsx"
            output.seek(0)
            return output, output_filename

        except Exception as e:
            # Capture any general exceptions and log them
            print(f"Error during PDF to Excel conversion: {e}")
            raise ValueError("Failed to convert PDF to Excel. Please check if the PDF contains valid tables.")

    elif format == 'pptx':
        try:
            prs = Presentation()
            reader = PdfReader(file)
            for page_num, page in enumerate(reader.pages):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                text_frame = text_box.text_frame
                text_frame.text = page.extract_text()
            prs.save(output)
            output_filename = "converted_file.pptx"
            output.seek(0)
            return output, output_filename

        except Exception as e:
            print(f"Error during PDF to PPTX conversion: {e}")
            raise ValueError("Failed to convert PDF to PPTX.")

    elif format in ['jpeg', 'png']:
        try:
            doc = fitz.open(stream=file.read(), filetype="pdf")
            img_format = format.upper()

            if len(doc) > 1:
                with zipfile.ZipFile(output, mode='w', compression=zipfile.ZIP_DEFLATED) as img_zip:
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)
                        pix = page.get_pixmap()
                        img_bytes = pix.tobytes(img_format)
                        img_filename = f'page_{page_num + 1}.{format}'
                        img_zip.writestr(img_filename, img_bytes)

                output_filename = "converted_pages.zip"
                output.seek(0)
                return output, output_filename

            else:
                pix = doc.load_page(0).get_pixmap()
                img_bytes = pix.tobytes(img_format)
                output.write(img_bytes)
                output_filename = f"converted_page.{format}"
                output.seek(0)
                return output, output_filename

        except fitz.FileDataError as e:
            print(f"Error: {e}")
            raise ValueError("Failed to process the PDF file. Please check if it's a valid PDF.")

    output.seek(0)
    return output, None


def html_to_pdf(html_content):
    from xhtml2pdf import pisa
    output = io.BytesIO()
    pisa.CreatePDF(html_content, dest=output)
    output.seek(0)
    return output
