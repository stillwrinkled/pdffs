import io
import os
import zipfile
from pdf2docx import Converter as PDF2DOCXConverter
import fitz  # PyMuPDF for converting PDF to images (JPEG/PNG)
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import camelot  # For extracting tables from PDF
import tabula  # Alternative for extracting tables using Java
import logging
from PIL import Image
import tempfile

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def convert_pdf(file, format):
    output = io.BytesIO()
    file.seek(0)  # Reset file pointer to the start

    if format == 'docx':
        return convert_pdf_to_docx(file)
    elif format == 'xlsx':
        return convert_pdf_to_excel(file)
    elif format == 'pptx':
        return convert_pdf_to_pptx(file)
    elif format in ['jpeg', 'png']:
        return convert_pdf_to_image(file, format)
    elif format == 'dcx':
        return convert_pdf_to_dcx(file)

    output.seek(0)
    return output, None


def convert_pdf_to_docx(file):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            temp_pdf.write(file.read())
            temp_pdf.close()
            output_file = io.BytesIO()
            converter = PDF2DOCXConverter(temp_pdf.name)
            converter.convert(output_file, start=0, end=None)
            converter.close()
            os.remove(temp_pdf.name)
            output_file.seek(0)
            output_filename = "converted_file.docx"
            logger.info("PDF to DOCX conversion successful")
            return output_file, output_filename

    except Exception as e:
        logger.error(f"Error during PDF to DOCX conversion: {e}")
        raise ValueError("Failed to convert PDF to DOCX.")


def convert_pdf_to_excel(file):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            temp_pdf.write(file.read())
            temp_pdf.close()

            try:
                logger.info("Attempting table extraction with Camelot")
                tables = camelot.read_pdf(temp_pdf.name, pages='all', flavor='stream')

                if not tables or len(tables) == 0:
                    raise ValueError("No tables found in the PDF.")

                wb = Workbook()
                ws = wb.active

                for table_num, table in enumerate(tables):
                    for i, row in enumerate(table.df.itertuples(), 1):
                        for j, value in enumerate(row[1:], 1):
                            ws.cell(row=i, column=j, value=value)

                output = io.BytesIO()
                wb.save(output)
                output_filename = "converted_file.xlsx"
                output.seek(0)
                os.remove(temp_pdf.name)
                logger.info("PDF to Excel conversion successful using Camelot")
                return output, output_filename

            except Exception as camelot_err:
                logger.error(f"Camelot failed: {camelot_err}")
                logger.info("Camelot failed, trying Tabula")
                try:
                    tables = tabula.read_pdf(temp_pdf.name, pages='all', multiple_tables=True)
                    if not tables or len(tables) == 0:
                        raise ValueError("No tables found using Tabula.")

                    wb = Workbook()
                    ws = wb.active

                    for table in tables:
                        for row in table.itertuples():
                            ws.append(row)

                    output = io.BytesIO()
                    wb.save(output)
                    output_filename = "converted_file.xlsx"
                    output.seek(0)
                    os.remove(temp_pdf.name)
                    logger.info("PDF to Excel conversion successful using Tabula")
                    return output, output_filename

                except Exception as tabula_err:
                    logger.error(f"Tabula failed: {tabula_err}")
                    raise ValueError("Failed to extract tables with both Camelot and Tabula.")

    except Exception as e:
        logger.error(f"Error during PDF to Excel conversion: {e}")
        raise ValueError("Failed to convert PDF to Excel. Please check if the PDF contains valid tables.")


def convert_pdf_to_pptx(file):
    try:
        prs = Presentation()
        doc = fitz.open(stream=file.read(), filetype="pdf")

        for page_num, page in enumerate(doc.pages(), start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Extract text instances
            text_instances = page.get_text("dict")["blocks"]

            # Existing text extraction and positioning logic
            # ...

            # Extract images and place on slide
            for img_index, img in enumerate(page.get_images(full=True), start=1):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    img_extension = base_image["ext"]

                    # Try to open the image to verify it's valid
                    image = Image.open(io.BytesIO(image_bytes))
                    img_file = f"image_{page_num}_{img_index}.{img_extension}"

                    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{img_extension}") as temp_img:
                        temp_img.write(image_bytes)
                        temp_img_path = temp_img.name
                        slide.shapes.add_picture(temp_img_path, Inches(0.5), Inches(2), width=Inches(4))

                except Exception as img_error:
                    logger.error(f"Error processing image on page {page_num}, image index {img_index}: {img_error}")
                    continue  # Skip this image and continue with the rest

        output = io.BytesIO()
        prs.save(output)
        output_filename = "converted_file.pptx"
        output.seek(0)
        logger.info("PDF to PPTX conversion successful")
        return output, output_filename

    except Exception as e:
        logger.error(f"Error during PDF to PPTX conversion: {e}")
        raise ValueError("Failed to convert PDF to PPTX.")


def convert_pdf_to_image(file, format):
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        img_format = format.upper()

        output = io.BytesIO()

        if len(doc) > 1:
            with zipfile.ZipFile(output, mode='w', compression=zipfile.ZIP_DEFLATED) as img_zip:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap()
                    img_bytes = pix.tobytes(img_format)
                    img_filename = f'page_{page_num + 1}.{format}'
                    img_zip.writestr(img_filename, img_bytes)

            output_filename = "converted_pages.zip"
            output.seek(0)
            logger.info("PDF to multi-page image (ZIP) conversion successful")
            return output, output_filename

        else:
            pix = doc.load_page(0).get_pixmap()
            img_bytes = pix.tobytes(img_format)
            output.write(img_bytes)
            output_filename = f"converted_page.{format}"
            output.seek(0)
            logger.info("PDF to single image conversion successful")
            return output, output_filename

    except fitz.FileDataError as e:
        logger.error(f"Error during image conversion: {e}")
        raise ValueError("Failed to process the PDF file. Please check if it's a valid PDF.")


def convert_pdf_to_dcx(file):
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        output = io.BytesIO()

        if len(doc) > 1:
            with zipfile.ZipFile(output, mode='w', compression=zipfile.ZIP_DEFLATED) as dcx_zip:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap()
                    img_bytes = pix.tobytes("DCX")
                    img_filename = f'page_{page_num + 1}.dcx'
                    dcx_zip.writestr(img_filename, img_bytes)

            output_filename = "converted_pages.dcx.zip"
            output.seek(0)
            logger.info("PDF to DCX conversion successful")
            return output, output_filename

        else:
            pix = doc.load_page(0).get_pixmap()
            img_bytes = pix.tobytes("DCX")
            output.write(img_bytes)
            output_filename = "converted_page.dcx"
            output.seek(0)
            logger.info("PDF to single-page DCX conversion successful")
            return output, output_filename

    except fitz.FileDataError as e:
        logger.error(f"Error during DCX conversion: {e}")
        raise ValueError("Failed to process the PDF file for DCX conversion.")
