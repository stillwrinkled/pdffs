import io
import os
import zipfile
import fitz  # PyMuPDF for PDF handling (fonts, images, etc.)
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pdf2docx import Converter as PDF2DOCXConverter
import camelot
import tabula
from PyPDF2 import PdfReader
import tempfile
import logging
from io import BytesIO

from PIL import Image  # For image processing

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def convert_pdf(file, format):
    output = io.BytesIO()
    file.seek(0)  # Reset file pointer to the start

    if format == 'docx':
        return convert_pdf_to_docx(file)
    elif format == 'xlsx':
        return convert_pdf_to_excel(file)
    elif format == 'pptx':
        return convert_pdf_to_pptx(file)
    elif format in ['jpeg', 'png']:
        return convert_pdf_to_image(file, format)

    output.seek(0)
    return output, None


def convert_pdf_to_docx(file):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            temp_pdf.write(file.read())
            temp_pdf.close()
            output_file = io.BytesIO()
            converter = PDF2DOCXConverter(temp_pdf.name)
            converter.convert(output_file, start=0, end=None)
            converter.close()
            os.remove(temp_pdf.name)
            output_file.seek(0)
            output_filename = "converted_file.docx"
            logger.info("PDF to DOCX conversion successful")
            return output_file, output_filename

    except Exception as e:
        logger.error(f"Error during PDF to DOCX conversion: {e}")
        raise ValueError("Failed to convert PDF to DOCX.")


def convert_pdf_to_excel(file):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            temp_pdf.write(file.read())
            temp_pdf.close()

            try:
                logger.info("Attempting table extraction with Camelot")
                tables = camelot.read_pdf(temp_pdf.name, pages='all', flavor='stream')

                if not tables or len(tables) == 0:
                    raise ValueError("No tables found in the PDF.")

                wb = Workbook()
                ws = wb.active

                for table_num, table in enumerate(tables):
                    for i, row in enumerate(table.df.itertuples(), 1):
                        for j, value in enumerate(row[1:], 1):
                            ws.cell(row=i, column=j, value=value)

                output = io.BytesIO()
                wb.save(output)
                output_filename = "converted_file.xlsx"
                output.seek(0)
                os.remove(temp_pdf.name)
                logger.info("PDF to Excel conversion successful using Camelot")
                return output, output_filename

            except Exception as camelot_err:
                logger.error(f"Camelot failed: {camelot_err}")
                logger.info("Camelot failed, trying Tabula")
                try:
                    tables = tabula.read_pdf(temp_pdf.name, pages='all', multiple_tables=True)
                    if not tables or len(tables) == 0:
                        raise ValueError("No tables found using Tabula.")

                    wb = Workbook()
                    ws = wb.active

                    for table in tables:
                        for row in table.itertuples():
                            ws.append(row)

                    output = io.BytesIO()
                    wb.save(output)
                    output_filename = "converted_file.xlsx"
                    output.seek(0)
                    os.remove(temp_pdf.name)
                    logger.info("PDF to Excel conversion successful using Tabula")
                    return output, output_filename

                except Exception as tabula_err:
                    logger.error(f"Tabula failed: {tabula_err}")
                    raise ValueError("Failed to extract tables with both Camelot and Tabula.")

    except Exception as e:
        logger.error(f"Error during PDF to Excel conversion: {e}")
        raise ValueError("Failed to convert PDF to Excel. Please check if the PDF contains valid tables.")


def convert_pdf_to_pptx(file):
    # Convert Flask's FileStorage to bytes
    file_bytes = file.read()

    # Open the PDF using PyMuPDF with bytes stream
    pdf_document = fitz.open(stream=file_bytes, filetype="pdf")

    # Create a new presentation object
    prs = Presentation()

    for page_num in range(pdf_document.page_count):
        # Extract the page
        page = pdf_document[page_num]

        # Create a new slide
        slide_layout = prs.slide_layouts[5]  # Use blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Extract text and images from the page
        text_instances = page.get_text("dict")
        image_instances = page.get_images(full=True)

        # Adding text to the slide
        for block in text_instances["blocks"]:
            if "lines" in block:
                for line in block["lines"]:
                    for span in line["spans"]:
                        text = span["text"].strip()
                        if text:
                            # Add a textbox with the text to the slide
                            left = Inches(0.5)  # Positioning of text box
                            top = Inches(1 + page_num)  # Offset based on page
                            width = Inches(9)
                            height = Inches(0.5)
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            p = text_frame.add_paragraph()
                            run = p.add_run()
                            run.text = text
                            run.font.size = Pt(12)
                            # Optional: Add some font styling (color, size)
                            if "color" in span:
                                run.font.color.rgb = RGBColor(0, 0, 0)  # Black color
                            if "size" in span:
                                run.font.size = Pt(span["size"])

        # Adding images to the slide
        for img in image_instances:
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]

            # Save the image to a file-like object
            image_stream = BytesIO(image_bytes)

            # Add the image to the slide
            left = Inches(0.5)
            top = Inches(1 + page_num)
            slide.shapes.add_picture(image_stream, left, top, width=Inches(3), height=Inches(2))

    # Save the presentation as a file-like object
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)

    # Reset the buffer's position to the beginning
    pptx_buffer.seek(0)

    # Return the PPTX file content and filename
    return pptx_buffer, "converted_output.pptx"


def convert_pdf_to_image(file, format):
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        img_format = format.upper()

        output = io.BytesIO()

        if len(doc) > 1:
            with zipfile.ZipFile(output, mode='w', compression=zipfile.ZIP_DEFLATED) as img_zip:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap()
                    img_bytes = pix.tobytes(img_format)
                    img_filename = f'page_{page_num + 1}.{format}'
                    img_zip.writestr(img_filename, img_bytes)

            output_filename = "converted_pages.zip"
            output.seek(0)
            logger.info("PDF to multi-page image (ZIP) conversion successful")
            return output, output_filename

        else:
            pix = doc.load_page(0).get_pixmap()
            img_bytes = pix.tobytes(img_format)
            output.write(img_bytes)
            output_filename = f"converted_page.{format}"
            output.seek(0)
            logger.info("PDF to single image conversion successful")
            return output, output_filename

    except fitz.FileDataError as e:
        logger.error(f"Error during image conversion: {e}")
        raise ValueError("Failed to process the PDF file. Please check if it's a valid PDF.")
